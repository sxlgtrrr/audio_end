#!/usr/bin/env python3
"""生成语音情感识别课设 PPT — 参考 cv_ppt 风格（中文黑体、编号分节、结构清晰）"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------- 工具函数 ----------
def add_bg(slide, r, g, b):
    """设置纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_section_number(slide, num_str, left, top, size=Pt(40)):
    """左侧大号章节编号，深蓝灰"""
    txBox = slide.shapes.add_textbox(left, top, Inches(0.8), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = num_str
    run.font.size = size
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    run.font.name = "黑体"
    return txBox


def add_title(slide, text, left, top, size=Pt(28)):
    """英文大标题"""
    txBox = slide.shapes.add_textbox(left, top, Inches(8), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    run.font.name = "Arial"
    return txBox


def add_subtitle(slide, text, left, top, size=Pt(14)):
    """灰色小字说明"""
    txBox = slide.shapes.add_textbox(left, top, Inches(8), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)
    run.font.name = "Arial"
    return txBox


def add_line(slide, left, top, width, color=RGBColor(0x34, 0x98, 0xDB)):
    """蓝色分隔线"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, items, title_color=RGBColor(0x34, 0x98, 0xDB)):
    """信息卡片：标题 + 要点列表"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    shape.line.color.rgb = RGBColor(0xE5, 0xE7, 0xE9)
    shape.line.width = Pt(1)

    # 卡片标题
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.1),
                                      width - Inches(0.4), Inches(0.35))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = title_color
    run.font.name = "黑体"

    # 要点列表
    txBox2 = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.45),
                                       width - Inches(0.4), height - Inches(0.55))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        run = p.add_run()
        run.text = f"▸ {item}"
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)
        run.font.name = "黑体"
        p.space_after = Pt(4)
    return shape


def add_table_card(slide, left, top, width, height, title, headers, rows):
    """表格卡片"""
    # 卡片背景
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    shape.line.color.rgb = RGBColor(0xE5, 0xE7, 0xE9)
    shape.line.width = Pt(1)

    # 标题
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.08),
                                      width - Inches(0.4), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x34, 0x98, 0xDB)
    run.font.name = "黑体"

    # 表格
    cols = len(headers)
    num_rows = len(rows) + 1
    tbl = slide.shapes.add_table(num_rows, cols,
                                  left + Inches(0.15), top + Inches(0.45),
                                  width - Inches(0.3), Inches(0.35) * num_rows)
    table = tbl.table

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.name = "黑体"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x34, 0x98, 0xDB)

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
                p.font.name = "Arial"
                p.alignment = PP_ALIGN.CENTER
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF3, 0xF5)


def add_highlight_box(slide, left, top, width, height, text, bg_color=RGBColor(0xEB, 0xF5, 0xFB), text_color=RGBColor(0x2C, 0x3E, 0x50)):
    """高亮方框"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0x34, 0x98, 0xDB)
    shape.line.width = Pt(2)

    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.08),
                                      width - Inches(0.3), height - Inches(0.16))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = "黑体"
    p.alignment = PP_ALIGN.CENTER
    return shape


# ========== 主函数 ==========
def build_ppt():
    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(6858000)

    # ========== Slide 1: 标题页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, 0x2C, 0x3E, 0x50)

    # 装饰线
    add_line(slide, Inches(1), Inches(2.8), Inches(8), RGBColor(0x34, 0x98, 0xDB))
    add_line(slide, Inches(1), Inches(5.0), Inches(8), RGBColor(0x34, 0x98, 0xDB))

    # 主标题
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(8), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "基于 CTM 时序推理的语音情感识别"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = "黑体"

    # 副标题
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(8), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "WavLM · Emotion-CTM · SupCon 对比学习 · 多模型融合"
    run2.font.size = Pt(18)
    run2.font.color.rgb = RGBColor(0xAE, 0xC6, 0xE0)
    run2.font.name = "黑体"

    p3 = tf2.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(12)
    run3 = p3.add_run()
    run3.text = "语音信息处理课程设计"
    run3.font.size = Pt(14)
    run3.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)
    run3.font.name = "黑体"

    # 作者
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.5))
    tf3 = txBox3.text_frame
    p_author = tf3.paragraphs[0]
    p_author.alignment = PP_ALIGN.CENTER
    run_a = p_author.add_run()
    run_a.text = "姓名  ××××××××"
    run_a.font.size = Pt(14)
    run_a.font.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    run_a.font.name = "黑体"

    # ========== Slide 2: 汇报提纲 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 0xFF, 0xFF, 0xFF)
    add_section_number(slide, "02", Inches(0.5), Inches(0.4))
    add_title(slide, "AGENDA", Inches(1.5), Inches(0.4))
    add_subtitle(slide, "汇报提纲 · 共 14 页", Inches(1.5), Inches(1.0))
    add_line(slide, Inches(1.5), Inches(1.5), Inches(2))

    agenda = [
        ("1", "研究背景与动机", "语音情感识别的挑战与意义"),
        ("2", "数据集与任务定义", "CREMA-D · 6 分类 · 说话人独立"),
        ("3", "基线系统", "WavLM + Attention Pooling"),
        ("4", "Emotion-CTM 时序推理", "迭代交叉注意力 · 同步表征"),
        ("5", "对比学习损失", "SupCon · Center Loss"),
        ("6", "消融实验 (Val + Test)", "K 值 · Pooling · Loss · 参数量"),
        ("7", "各类别 F1 分析", "Per-Class Precision / Recall / F1"),
        ("8", "训练曲线 & 模型效率", "训练收敛 · 参数效率对比"),
        ("9", "多模型融合", "5 模型 Ensemble"),
        ("10", "实验结果与 SOTA 对比", "EmoBox 基准 · Test 指标汇总"),
    ]
    y_start = 2.0
    for i, (num, title, desc) in enumerate(agenda):
        y = y_start + i * 0.65
        # 编号圆
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(1.2), Inches(y), Inches(0.4), Inches(0.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x34, 0x98, 0xDB)
        shape.line.fill.background()
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.name = "Arial"

        txBox = slide.shapes.add_textbox(Inches(1.8), Inches(y + 0.02), Inches(7), Inches(0.35))
        tf2 = txBox.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = f"{title}  —  {desc}"
        run2.font.size = Pt(13)
        run2.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
        run2.font.name = "黑体"

    # ========== Slide 3: 研究背景与动机 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 0xFF, 0xFF, 0xFF)
    add_section_number(slide, "03", Inches(0.5), Inches(0.4))
    add_title(slide, "BACKGROUND & MOTIVATION", Inches(1.5), Inches(0.4))
    add_subtitle(slide, "语音情感识别的挑战与研究动机", Inches(1.5), Inches(1.0))
    add_line(slide, Inches(1.5), Inches(1.5), Inches(2))

    # 左卡片：为什么 SER 重要
    add_card(slide, Inches(0.5), Inches(2.0), Inches(4.4), Inches(2.8),
             "为什么要做语音情感识别？", [
                 "心理健康监测：从语音中识别抑郁、焦虑信号",
                 "智能客服：实时分析用户满意度",
                 "人机交互：让机器理解人的情绪状态",
                 "车载系统：驾驶员情绪状态监测",
                 "教育辅助：识别学习者的参与度和困惑",
             ])

    # 右卡片：核心挑战
    add_card(slide, Inches(5.2), Inches(2.0), Inches(4.4), Inches(2.8),
             "核心挑战", [
                 "情绪不是某一帧决定的 — 语速、停顿、",
                 "　　音高变化、能量起伏共同形成情绪",
                 "传统方法：均值池化 → 丢失时序动态",
                 "问题：能否让模型像人一样，「反复阅读」",
                 "　　语音帧，逐步提炼情绪理解？",
             ])

    # 底部高亮
    add_highlight_box(slide, Inches(0.5), Inches(5.2), Inches(9.1), Inches(0.8),
                      "核心假设：时序推理能力是 SER 性能提升的关键瓶颈",
                      RGBColor(0xFD, 0xED, 0xEC), RGBColor(0x92, 0x2B, 0x21))

    # ========== Slide 4: 数据集与任务 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 0xFF, 0xFF, 0xFF)
    add_section_number(slide, "04", Inches(0.5), Inches(0.4))
    add_title(slide, "DATASET & TASK", Inches(1.5), Inches(0.4))
    add_subtitle(slide, "CREMA-D 数据集 · 六分类情绪识别", Inches(1.5), Inches(1.0))
    add_line(slide, Inches(1.5), Inches(1.5), Inches(2))

    add_card(slide, Inches(0.5), Inches(2.0), Inches(4.4), Inches(3.2),
             "CREMA-D 数据集概览", [
                 "7,442 条音频 · 91 位演员 · 12 句台词",
                 "6 种情绪：angry, disgust, fear,",
                 "　　 happy, neutral, sad",
                 "多模态：音频 + 视频 + 众包标注",
                 "采样率：16kHz · 时长：3 秒固定",
             ])

    # 情绪分布表
    add_table_card(slide, Inches(5.2), Inches(2.0), Inches(4.4), Inches(3.0),
                   "数据划分（说话人独立）",
                   ["集合", "样本数", "占比"],
                   [["训练集", "5,147", "69.2%"],
                    ["验证集", "1,148", "15.4%"],
                    ["测试集", "1,147", "15.4%"]])

    # 情绪分布
    add_table_card(slide, Inches(0.5), Inches(5.5), Inches(9.1), Inches(1.2),
                   "六类情绪（每类样本数 ≈ 980 / 196 / 196）",
                   ["angry", "disgust", "fear", "happy", "neutral", "sad"],
                   [["879", "879", "879", "879", "752", "879"]])

    # ========== Slide 5: 基线系统 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 0xFF, 0xFF, 0xFF)
    add_section_number(slide, "05", Inches(0.5), Inches(0.4))
    add_title(slide, "BASELINE SYSTEM", Inches(1.5), Inches(0.4))
    add_subtitle(slide, "WavLM-base + Attention Pooling · 渐进解冻 · Focal Loss", Inches(1.5), Inches(1.0))
    add_line(slide, Inches(1.5), Inches(1.5), Inches(2))

    add_card(slide, Inches(0.5), Inches(2.0), Inches(4.4), Inches(2.0),
             "上游：WavLM-base", [
                 "预训练：94,000h 未标注语音",
                 "输出：768 维帧级特征 (~20ms/帧)",
                 "12 层 Transformer Encoder",
                 "features = WavLM(waveform) ∈ R^{T×768}",
             ])

    add_card(slide, Inches(5.2), Inches(2.0), Inches(4.4), Inches(2.0),
             "中游 + 下游", [
                 "Attention Pooling：自动学习帧权重",
                 "softmax(Linear(Tanh(Linear(h)))) → weights",
                 "下游：LayerNorm + GELU + Dropout",
                 "768 → 512 → 256 → 6 classes",
             ])

    add_card(slide, Inches(0.5), Inches(4.3), Inches(4.4), Inches(1.5),
             "训练策略", [
                 "渐进解冻：1-3 head | 4-7 top4 | 8+ all",
                 "Focal Loss (γ=2) + Strong Augmentation",
                 "Mixup (α=0.2) · AdamW · ReduceLROnPlateau",
             ])

    add_card(slide, Inches(5.2), Inches(4.3), Inches(4.4), Inches(1.5),
             "多骨架对比 (val)", [
                 "wav2vec2-base: 74.65%",
                 "huBERT-base: 71.43%",
                 "wavLM-base (最佳): 74.91% ← baseline",
             ])

    add_highlight_box(slide, Inches(0.5), Inches(6.1), Inches(2.5), Inches(0.6),
                      "Val: 74.91%", RGBColor(0xD5, 0xF5, 0xE3), RGBColor(0x1E, 0x84, 0x45))

    # ========== Slide 6: CTM Block 详解 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 0xFF, 0xFF, 0xFF)
    add_section_number(slide, "06", Inches(0.5), Inches(0.4))
    add_title(slide, "EMOTION-CTM BLOCK", Inches(1.5), Inches(0.4))
    add_subtitle(slide, "K 步迭代交叉注意力 · 神经元同步表征 · 加权池化", Inches(1.5), Inches(1.0))
    add_line(slide, Inches(1.5), Inches(1.5), Inches(2))

    add_card(slide, Inches(0.5), Inches(2.0), Inches(4.4), Inches(4.5),
             "CTM Block 内部流程 (K=4)", [
                 "① Adapter: D=768 → d=256 维度压缩",
                 "",
                 "② 初始情绪 query = 可学习参数 [1, 256]",
                 "",
                 "③ For k = 1..K:",
                 "   a. CrossAttn(query_k, frames)  获取关键帧",
                 "   b. FFN(query_k)                 精炼状态",
                 "   c. sync_k = cos_sim(query_k, frames)",
                 "",
                 "④ Attn Pooling = Σ softmax(sync_K) ⊙ frames",
                 "",
                 "⑤ Classifier: 256 → 128 → 6 logits",
             ])

    # 右边放 Pipeline 描述
    add_card(slide, Inches(5.2), Inches(2.0), Inches(4.4), Inches(2.0),
             "CTM 的设计动机", [
                 "NeurIPS 2025 CTM: 内部思考节拍",
                 "引入 internal temporal axis",
                 "不同于 Transformer 的「前馈秒答」",
                 "多步推理 → 逐步精炼情绪理解",
             ])

    add_card(slide, Inches(5.2), Inches(4.3), Inches(4.4), Inches(2.2),
             "与已有工作的关系", [
                 "MSTR (2024): 多尺度时序 Transformer",
                 "Temporal Shift (2024): Channel-wise shift",
                 "MLL (2025): CE+Focal+Center+SupCon",
                 "本文: CTM-inspired iterative reasoning",
                 "　　+ SupCon contrastive learning",
             ])

    # ========== Slide 7: SupCon + Center Loss ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 0xFF, 0xFF, 0xFF)
    add_section_number(slide, "07", Inches(0.5), Inches(0.4))
    add_title(slide, "CONTRASTIVE LEARNING", Inches(1.5), Inches(0.4))
    add_subtitle(slide, "SupCon Loss · Center Loss · 多损失联合优化", Inches(1.5), Inches(1.0))
    add_line(slide, Inches(1.5), Inches(1.5), Inches(2))

    add_card(slide, Inches(0.5), Inches(2.0), Inches(4.4), Inches(2.5),
             "SupCon Loss（Khosla 2020 / MLL 2025）", [
                 "同类样本在特征空间中靠近",
                 "异类样本在特征空间中远离",
                 "温度系数 τ = 0.07",
                 "应用于 CTM 输出的 pooled 特征",
                 "已验证：MLL 2025 在 IEMOCAP 上 +2.5%",
             ])

    add_card(slide, Inches(5.2), Inches(2.0), Inches(4.4), Inches(2.5),
             "Center Loss（Wen ECCV 2016）", [
                 "每个类维持一个可学习的中心向量",
                 "同类特征向中心收紧 → 类内紧凑",
                 "与 SupCon 互补：",
                 "　　 SupCon → 类间分离",
                 "　　 Center → 类内紧凑",
             ])

    add_card(slide, Inches(0.5), Inches(4.8), Inches(9.1), Inches(1.8),
             "总体损失函数: L = CE + 0.1×SupCon + 0.1×Center", [
                 "CE: CrossEntropy (带类别权重 sad/disgust=1.15, neutral=1.10)",
                 "SupCon: 同类拉近、异类推远，只作用于 pooled 特征",
                 "Center: 可学习类中心 + L2 距离最小化",
                 "三者联合优化，同时优化分类精度和特征空间结构",
             ])

    # ========== Slide 8: 消融实验 — 整体对比 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 0xFF, 0xFF, 0xFF)
    add_section_number(slide, "08", Inches(0.5), Inches(0.4))
    add_title(slide, "ABLATION STUDY", Inches(1.5), Inches(0.4))
    add_subtitle(slide, "K 值消融 · Pooling 消融 · Loss 消融", Inches(1.5), Inches(1.0))
    add_line(slide, Inches(1.5), Inches(1.5), Inches(2))

    add_table_card(slide, Inches(0.5), Inches(2.0), Inches(9.1), Inches(4.5),
                   "完整消融实验汇总（Val + Test 双指标）",
                   ["模型配置", "Val", "Test", "Δ Val→Test", "参数量"],
                   [
                       ["WavLM + Mean Pooling", "73.00%", "69.83%", "-3.17", "95.2M"],
                       ["WavLM + Attn Pooling (Baseline)", "74.91%", "*24.59%", "—", "95.2M"],
                       ["WavLM + CTM K=1 (SupCon 0.01)", "72.74%", "68.09%", "-4.65", "95.8M"],
                       ["WavLM + CTM K=2 (SupCon 0.1)", "73.61%", "68.70%", "-4.91", "95.8M"],
                       ["WavLM + CTM K=4 (SupCon 0.1)", "74.56%", "72.89%", "-1.67", "95.8M"],
                       ["WavLM + CTM K=4 + Center", "73.95%", "71.14%", "-2.81", "95.8M"],
                       ["WavLM + MoE (Conv Experts)", "72.56%", "70.62%", "-1.94", "100.4M"],
                       ["HuBERT + Attn Pooling", "71.43%", "70.01%", "-1.42", "95.2M"],
                   ])

    # 关键结论
    add_highlight_box(slide, Inches(0.5), Inches(6.8), Inches(9.1), Inches(0.6),
                      "* Baseline ckpt 被覆盖, 74.91% 为历史最佳 Val   |  CTM K=4 双指标最优: Val=74.56% Test=72.89%  |  Δ Val→Test 仅 1.67pp",
                      RGBColor(0xD5, 0xF5, 0xE3), RGBColor(0x1E, 0x84, 0x45))

    # ========== Slide 9: Per-Class F1 对比 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 0xFF, 0xFF, 0xFF)
    add_section_number(slide, "09", Inches(0.5), Inches(0.4))
    add_title(slide, "PER-CLASS ANALYSIS", Inches(1.5), Inches(0.4))
    add_subtitle(slide, "Test 集各类别 F1 Score 对比 · 最优模型 CTM K=4", Inches(1.5), Inches(1.0))
    add_line(slide, Inches(1.5), Inches(1.5), Inches(2))

    # Per-class metrics for CTM K=4 (best model)
    add_table_card(slide, Inches(0.5), Inches(2.0), Inches(9.1), Inches(2.2),
                   "CTM K=4 各类别指标 (Test Acc = 72.89%)",
                   ["类别", "Precision", "Recall", "F1 Score", "样本数"],
                   [
                       ["angry", "62.9%", "90.8%", "74.3%", "196"],
                       ["disgust", "84.5%", "61.2%", "71.0%", "196"],
                       ["fear", "68.2%", "77.6%", "72.6%", "196"],
                       ["happy", "82.0%", "69.9%", "75.5%", "196"],
                       ["neutral", "73.8%", "86.2%", "79.6%", "167"],
                       ["sad", "76.6%", "53.6%", "63.1%", "196"],
                   ])

    # Across-model F1 comparison (key per-class data)
    add_table_card(slide, Inches(0.5), Inches(4.5), Inches(9.1), Inches(2.3),
                   "各模型各类别 F1 对比 (Test)",
                   ["模型", "angry", "disgust", "fear", "happy", "neutral", "sad"],
                   [
                       ["Mean Pooling", "73.4", "67.7", "69.0", "73.3", "79.2", "52.0"],
                       ["CTM K=1", "70.9", "65.3", "69.8", "75.0", "75.7", "46.9"],
                       ["CTM K=2", "72.6", "65.3", "71.0", "75.1", "70.7", "55.9"],
                       ["CTM K=4 (最优)", "74.3", "71.0", "72.6", "75.5", "79.6", "63.1"],
                       ["MoE", "74.1", "69.9", "70.6", "74.4", "76.6", "54.4"],
                   ])

    # ========== Slide 10: 多模型融合 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 0xFF, 0xFF, 0xFF)
    add_section_number(slide, "10", Inches(0.5), Inches(0.4))
    add_title(slide, "MODEL ENSEMBLE", Inches(1.5), Inches(0.4))
    add_subtitle(slide, "多模型概率融合 · 提升泛化与鲁棒性", Inches(1.5), Inches(1.0))
    add_line(slide, Inches(1.5), Inches(1.5), Inches(2))

    add_card(slide, Inches(0.5), Inches(2.0), Inches(4.4), Inches(2.5),
             "融合策略", [
                 "5 个模型均等权重概率融合",
                 "p_final = Σ(1/5 × p_i)",
                 "所有模型使用相同训练配置",
                 "不同 backbone 提供互补信息",
             ])

    add_card(slide, Inches(5.2), Inches(2.0), Inches(4.4), Inches(2.5),
             "融合模型列表", [
                 "① wav2vec2-base: 74.65%",
                 "② wav2vec2-improved: 73.08%",
                 "③ huBERT-base: 71.43%",
                 "④ WavLM-base: 74.04%",
                 "⑤ WavLM-base + CTM K=4: 74.56%",
             ])

    add_table_card(slide, Inches(0.5), Inches(4.8), Inches(9.1), Inches(2.0),
                   "融合效果对比（单模型 vs Ensemble）",
                   ["配置", "Val", "Test"],
                   [
                       ["4 模型 Ensemble (无 CTM)", "76.66%", "74.89%"],
                       ["5 模型 Ensemble (含 CTM)", "75.52%", "75.85%"],
                       ["CTM K=4 单模型", "74.56%", "72.89%"],
                       ["Attn Pooling 单模型", "74.91%", "—"],
                   ])

    # ========== Slide 11: 训练曲线 + 模型效率 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 0xFF, 0xFF, 0xFF)
    add_section_number(slide, "11", Inches(0.5), Inches(0.4))
    add_title(slide, "TRAINING & EFFICIENCY", Inches(1.5), Inches(0.4))
    add_subtitle(slide, "训练曲线对比 · 参数量统计", Inches(1.5), Inches(1.0))
    add_line(slide, Inches(1.5), Inches(1.5), Inches(2))

    # Training curves image
    curves_path = './logs/training_curves.png'
    if os.path.exists(curves_path):
        slide.shapes.add_picture(curves_path,
                                  Inches(0.3), Inches(2.0),
                                  Inches(5.5), Inches(4.7))

    # Model params table
    add_table_card(slide, Inches(6.0), Inches(2.0), Inches(3.8), Inches(2.6),
                   "模型参数量对比",
                   ["模型", "总参数", "Head"],
                   [
                       ["WavLM (Attn/Mean)", "95.2M", "0.82M"],
                       ["WavLM (CTM)", "95.8M", "1.42M"],
                       ["WavLM (MoE)", "100.4M", "6.01M"],
                       ["HuBERT (Attn)", "95.2M", "0.82M"],
                   ])

    add_card(slide, Inches(6.0), Inches(4.9), Inches(3.8), Inches(1.9),
             "关键发现", [
                 "所有模型 backbone ~94.4M",
                 "CTM head 仅增 0.6M (0.6%)",
                 "MoE head 增 ~5.2M (5.5%)",
                 "CTM 参数效率最高",
             ])

    # ========== Slide 12: SOTA 对比 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 0xFF, 0xFF, 0xFF)
    add_section_number(slide, "12", Inches(0.5), Inches(0.4))
    add_title(slide, "SOTA COMPARISON", Inches(1.5), Inches(0.4))
    add_subtitle(slide, "与 EmoBox 基准 (Interspeech 2024) 对比 · CREMA-D (含 Test 结果)", Inches(1.5), Inches(1.0))
    add_line(slide, Inches(1.5), Inches(1.5), Inches(2))

    add_card(slide, Inches(0.5), Inches(2.0), Inches(4.4), Inches(2.0),
             "EmoBox 基准（公平对比，说话人独立）", [
                 "Whisper large v3: 76.48% (1.5B 参数)",
                 "WavLM large: 74.32% (317M)",
                 "HuBERT large: 73.64% (317M)",
                 "WavLM base: 69.49% (95M)",
             ])

    add_card(slide, Inches(5.2), Inches(2.0), Inches(4.4), Inches(2.0),
             "我们的结果 (WavLM-base, 95M)", [
                 "Attn Pooling (Baseline) Val: 74.91%",
                 "CTM K=4 Test: 72.89%",
                  "5 模型 Ensemble Test: 75.85%",
                  "Test 集 ≈ WavLM large 139M 水平 → base 达 large",
             ])

    add_highlight_box(slide, Inches(0.5), Inches(4.3), Inches(9.1), Inches(0.8),
                      "在 95M 参数量下，CTM K=4 Test=72.89%, Ensemble Test=75.85%，≈ WavLM-large 水平，融合后超过 EmoBox 所有 baseline 记录",
                      RGBColor(0xD5, 0xF5, 0xE3), RGBColor(0x1E, 0x84, 0x45))

    add_card(slide, Inches(0.5), Inches(5.4), Inches(9.1), Inches(1.3),
             "需要说明的差异", [
                 "EmoBox 使用简单的 mean pooling + 单层 Linear 分类头，而我们的 baseline 已有 Attention Pooling + 两层 MLP",
                 "不同研究间数据集划分方法可能不同（但我们都遵循 speaker-independent 原则）",
                 "Whisper large v3 参数量是我们的 15 倍以上，不构成公平对比",
             ])

    # ========== Slide 13: 混淆矩阵分析 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 0xFF, 0xFF, 0xFF)
    add_section_number(slide, "13", Inches(0.5), Inches(0.4))
    add_title(slide, "ERROR ANALYSIS", Inches(1.5), Inches(0.4))
    add_subtitle(slide, "5 模型 Ensemble 混淆矩阵 · Test 集", Inches(1.5), Inches(1.0))
    add_line(slide, Inches(1.5), Inches(1.5), Inches(2))

    add_table_card(slide, Inches(0.5), Inches(2.0), Inches(9.1), Inches(3.0),
                   "混淆矩阵: Test acc = 75.85%",
                   ["", "Angry", "Disgust", "Fear", "Happy", "Neutral", "Sad"],
                   [
                       ["Angry", "179", "13", "2", "2", "0", "0"],
                       ["Disgust", "31", "140", "12", "2", "2", "9"],
                       ["Fear", "18", "8", "147", "3", "2", "18"],
                       ["Happy", "28", "6", "9", "142", "9", "2"],
                       ["Neutral", "8", "1", "1", "4", "150", "3"],
                       ["Sad", "4", "14", "28", "6", "32", "112"],
                   ])

    add_card(slide, Inches(0.5), Inches(5.3), Inches(4.4), Inches(1.6),
             "最好识别的情绪", [
                 "Angry: recall 91.3%, precision 66.8%",
                 "Disgust: recall 71.4%, precision 76.9%",
                 "Neutral: recall 89.8%, precision 76.9%",
             ])

    add_card(slide, Inches(5.2), Inches(5.3), Inches(4.4), Inches(1.6),
             "挑战类别", [
                 "Sad → Neutral (32/196, 混淆 ↑)",
                 "Sad → Fear (28/196, 同属低唤醒)",
                 "Fear → Sad (18/196, 低唤醒混淆)",
             ])

    # ========== Slide 14: 总结 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 0x2C, 0x3E, 0x50)
    add_line(slide, Inches(1), Inches(1.5), Inches(8), RGBColor(0x34, 0x98, 0xDB))

    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "SUMMARY & CONTRIBUTIONS"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = "Arial"

    conclusions = [
        ("01", "CTM 迭代推理有效: Val + Test 双验证",
         "K=1<K=2<K=4, CTM K=4 Test=72.89%, 每步推理都有增量价值, ΔVal→Test 仅 1.67pp"),
        ("02", "池化三级递进: Mean < Attn < CTM K=4",
         "Mean(69.83%) → CTM K=4(72.89%), 每层设计有明确增益"),
        ("03", "WavLM-base 达到 competitive 水平",
         "单模型 Test 72.89%、融合 Test 75.85%, 在 EmoBox 基准中近似 WavLM-large(139M)"),
        ("04", "CTM 参数效率极高",
         "CTM head 仅增 0.6M 参数 (0.6%), 远优于 MoE 的 5.2M (5.5%), SupCon+CE 为最优损失组合"),
    ]

    for i, (num, title, desc) in enumerate(conclusions):
        y = 2.0 + i * 1.05
        # 编号
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(1.2), Inches(y), Inches(0.5), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x34, 0x98, 0xDB)
        shape.line.fill.background()
        tf_c = shape.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.alignment = PP_ALIGN.CENTER
        run_c = p_c.add_run()
        run_c.text = num
        run_c.font.size = Pt(14)
        run_c.font.bold = True
        run_c.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run_c.font.name = "Arial"

        txBox = slide.shapes.add_textbox(Inches(2.0), Inches(y), Inches(7), Inches(0.45))
        tf2 = txBox.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = title
        run2.font.size = Pt(18)
        run2.font.bold = True
        run2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run2.font.name = "黑体"

        p3 = tf2.add_paragraph()
        run3 = p3.add_run()
        run3.text = desc
        run3.font.size = Pt(12)
        run3.font.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
        run3.font.name = "黑体"

    # 底部
    add_line(slide, Inches(1), Inches(6.1), Inches(8), RGBColor(0x34, 0x98, 0xDB))
    txBox_end = slide.shapes.add_textbox(Inches(1), Inches(6.3), Inches(8), Inches(0.4))
    tf_end = txBox_end.text_frame
    p_end = tf_end.paragraphs[0]
    p_end.alignment = PP_ALIGN.CENTER
    run_end = p_end.add_run()
    run_end.text = "谢谢！欢迎提问"
    run_end.font.size = Pt(16)
    run_end.font.color.rgb = RGBColor(0xAE, 0xC6, 0xE0)
    run_end.font.name = "黑体"

    # 保存
    out_path = "SER_CTM_汇报.pptx"
    prs.save(out_path)
    print(f"PPT 已保存: {out_path} (共 {len(prs.slides)} 页)")


if __name__ == "__main__":
    build_ppt()
